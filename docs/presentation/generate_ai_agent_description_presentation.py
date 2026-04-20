from __future__ import annotations

from pathlib import Path
import sys
from textwrap import dedent

LOCAL_DEPS = Path(__file__).resolve().parent / "_pptx_deps"
if LOCAL_DEPS.exists():
    sys.path.insert(0, str(LOCAL_DEPS))

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Inches, Pt
except ModuleNotFoundError as exc:
    raise SystemExit(
        "Install python-pptx first: python -m pip install python-pptx -t docs\\presentation\\_pptx_deps"
    ) from exc


OUT_DIR = Path(__file__).resolve().parent
PPTX_PATH = OUT_DIR / "FinSense_AI_agent_description.pptx"
NOTES_PATH = OUT_DIR / "FinSense_AI_agent_description_speaker_notes.md"

WIDE_W = 13.333
WIDE_H = 7.5
FONT = "Arial"

INK = RGBColor(22, 33, 46)
MUTED = RGBColor(92, 108, 126)
LIGHT = RGBColor(248, 250, 252)
WHITE = RGBColor(255, 255, 255)
LINE = RGBColor(207, 217, 228)
TEAL = RGBColor(15, 148, 136)
TEAL_LIGHT = RGBColor(225, 247, 244)
BLUE = RGBColor(35, 99, 171)
BLUE_LIGHT = RGBColor(229, 240, 255)
AMBER = RGBColor(190, 111, 15)
AMBER_LIGHT = RGBColor(255, 242, 221)
GREEN = RGBColor(22, 130, 74)
GREEN_LIGHT = RGBColor(228, 246, 237)
RED = RGBColor(179, 53, 67)
RED_LIGHT = RGBColor(255, 234, 237)
PURPLE = RGBColor(105, 83, 168)
PURPLE_LIGHT = RGBColor(241, 238, 255)


def set_run(run, size: float, color=INK, bold: bool = False) -> None:
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold


def add_bg(slide, index: int | None = None) -> None:
    bg = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(WIDE_W), Inches(WIDE_H)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHT
    bg.line.fill.background()

    top = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(WIDE_W), Inches(0.12)
    )
    top.fill.solid()
    top.fill.fore_color.rgb = TEAL
    top.line.fill.background()

    side = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(12.82),
        Inches(0),
        Inches(0.51),
        Inches(WIDE_H),
    )
    side.fill.solid()
    side.fill.fore_color.rgb = BLUE_LIGHT
    side.line.fill.background()

    if index is not None:
        text(slide, 12.15, 6.92, 0.42, 0.25, str(index), 9, MUTED, PP_ALIGN.RIGHT)


def text(slide, x, y, w, h, value, size=16, color=INK, align=PP_ALIGN.LEFT, bold=False):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.text = value
    p.alignment = align
    p.space_after = Pt(2)
    set_run(p.runs[0], size, color, bold)
    return shape


def box(
    slide,
    x,
    y,
    w,
    h,
    items,
    fill=WHITE,
    line=LINE,
    align=PP_ALIGN.LEFT,
    radius=True,
):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.13)
    tf.margin_right = Inches(0.13)
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.08)
    for i, item in enumerate(items):
        value, size, color, bold = item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = value
        p.alignment = align
        p.space_after = Pt(3)
        set_run(p.runs[0], size, color, bold)
    return shape


def title(slide, heading: str, subtitle: str | None = None) -> None:
    text(slide, 0.58, 0.38, 11.5, 0.48, heading, 25, INK, bold=True)
    if subtitle:
        text(slide, 0.6, 0.9, 10.7, 0.34, subtitle, 12.5, MUTED)


def bullets(lines: list[str], size=13.5, color=INK):
    return [(f"• {line}", size, color, False) for line in lines]


def pill(slide, x, y, w, label, fill=BLUE_LIGHT, color=BLUE):
    return box(slide, x, y, w, 0.36, [(label, 11.5, color, True)], fill, fill, PP_ALIGN.CENTER)


def arrow(slide, x1, y1, x2, y2, color=TEAL):
    conn = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(1.9)
    return conn


def table(slide, x, y, widths, row_h, data, header_fill=INK, font_size=10.5):
    rows = len(data)
    cols = len(data[0])
    shape = slide.shapes.add_table(
        rows,
        cols,
        Inches(x),
        Inches(y),
        Inches(sum(widths)),
        Inches(row_h * rows),
    )
    native_table = shape.table

    for c, width in enumerate(widths):
        native_table.columns[c].width = Inches(width)
    for r in range(rows):
        native_table.rows[r].height = Inches(row_h)

    for r, row in enumerate(data):
        for c, value in enumerate(row):
            cell = native_table.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_fill if r == 0 else WHITE
            cell.margin_left = Inches(0.06)
            cell.margin_right = Inches(0.06)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)

            tf = cell.text_frame
            tf.clear()
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.text = value
            p.alignment = PP_ALIGN.CENTER
            p.space_after = Pt(0)
            set_run(p.runs[0], font_size, WHITE if r == 0 else INK, r == 0)

    return shape


def metric_card(slide, x, y, label, baseline, hybrid, color=TEAL):
    box(slide, x, y, 2.72, 1.0, [(label, 11.5, MUTED, False), (baseline, 20, RED, True), ("baseline", 9.5, MUTED, False)])
    arrow(slide, x + 2.85, y + 0.5, x + 3.28, y + 0.5, color)
    box(
        slide,
        x + 3.38,
        y,
        2.72,
        1.0,
        [(label, 11.5, MUTED, False), (hybrid, 20, GREEN, True), ("hybrid", 9.5, MUTED, False)],
        GREEN_LIGHT,
        GREEN,
    )


def slide1(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, 1)
    text(s, 0.72, 0.9, 10.6, 0.72, "Разработка и внедрение AI-агента", 33, INK, bold=True)
    text(s, 0.75, 1.58, 10.4, 0.58, "для оптимизации бизнес-процессов", 29, BLUE, bold=True)
    text(
        s,
        0.77,
        2.35,
        9.8,
        0.55,
        "Анализ существующих решений и разработка прототипа на базе Spring и Kotlin",
        17,
        MUTED,
    )
    box(s, 0.8, 3.45, 3.2, 1.0, [("Формат", 12, MUTED, False), ("магистерская диссертация", 19, INK, True)])
    box(s, 4.45, 3.45, 3.2, 1.0, [("Прототип", 12, MUTED, False), ("Spring + Kotlin", 20, TEAL, True)], TEAL_LIGHT, TEAL)
    box(s, 8.1, 3.45, 3.2, 1.0, [("Фокус", 12, MUTED, False), ("измеримый эффект процесса", 18, INK, True)])
    box(s, 1.0, 5.45, 2.75, 0.82, [("Бизнес-процесс", 17, INK, True)], BLUE_LIGHT, BLUE, PP_ALIGN.CENTER)
    arrow(s, 3.9, 5.86, 4.75, 5.86)
    box(s, 4.9, 5.45, 2.75, 0.82, [("AI-агент", 17, INK, True)], TEAL_LIGHT, TEAL, PP_ALIGN.CENTER)
    arrow(s, 7.8, 5.86, 8.65, 5.86)
    box(s, 8.8, 5.45, 2.75, 0.82, [("Измеримый эффект", 17, INK, True)], GREEN_LIGHT, GREEN, PP_ALIGN.CENTER)


def slide2(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, 2)
    title(s, "Актуальность и проблема исследования")
    box(
        s,
        0.7,
        1.42,
        5.8,
        4.35,
        [
            ("Почему нужен гибридный подход", 18, INK, True),
            *bullets(
                [
                    "бизнес-процессы включают типовые и неоднозначные операции",
                    "rule-based автоматизация эффективна только в заранее описанных сценариях",
                    "full LLM повышает гибкость, но увеличивает стоимость и задержки",
                    "практическая цель: улучшить измеримые параметры процесса",
                    "нужен баланс управляемости, адаптивности и реализуемости",
                ],
                13.2,
            ),
        ],
    )
    table(
        s,
        6.85,
        1.55,
        [1.55, 1.6, 1.65, 1.38],
        0.56,
        [
            ["Подход", "Управляемость", "Гибкость", "Стоимость"],
            ["Rule-based", "высокая", "низкая", "низкая"],
            ["Hybrid", "средняя / высокая", "средняя / высокая", "умеренная"],
            ["Full LLM", "низкая / средняя", "высокая", "высокая"],
        ],
        font_size=9.8,
    )
    box(
        s,
        7.05,
        4.45,
        5.72,
        1.18,
        [("Исследовательская задача", 14, BLUE, True), ("встроить LLM только в сложный сегмент процесса, сохранив быстрый и контролируемый основной путь", 15, INK, False)],
        BLUE_LIGHT,
        BLUE,
    )


def slide3(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, 3)
    title(s, "Объект, предмет, цель, задачи и методы")
    box(s, 0.7, 1.35, 3.8, 1.1, [("Объект", 17, BLUE, True), ("бизнес-процессы с классификацией и принятием решений при частичной неопределённости", 13.2, INK, False)], BLUE_LIGHT, BLUE)
    box(s, 4.8, 1.35, 3.8, 1.1, [("Предмет", 17, TEAL, True), ("архитектурные и программные подходы к LLM-агенту в event-driven процессе", 13.2, INK, False)], TEAL_LIGHT, TEAL)
    box(s, 8.9, 1.35, 3.35, 1.1, [("Методы", 17, AMBER, True), ("сравнительный анализ, проектирование, прототипирование, сценарное моделирование", 13, INK, False)], AMBER_LIGHT, AMBER)
    box(
        s,
        0.7,
        2.9,
        5.35,
        2.7,
        [
            ("Цель исследования", 18, INK, True),
            ("разработать и исследовать прототип AI-агента, встроенного в событийно-ориентированный бизнес-процесс", 15.5, INK, False),
            ("гибридный подход: rule-based логика + LLM fallback", 14.5, TEAL, True),
        ],
    )
    box(
        s,
        6.45,
        2.9,
        5.8,
        2.7,
        [
            ("Задачи", 18, INK, True),
            *bullets(
                [
                    "проанализировать подходы к AI-автоматизации",
                    "уточнить понятие AI-агента для backend-системы",
                    "разработать архитектуру прототипа",
                    "реализовать прототип на Spring + Kotlin",
                    "провести демонстрационную оценку",
                ],
                12.7,
            ),
        ],
    )


def slide4(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, 4)
    title(s, "AI-агент в контексте работы")
    box(
        s,
        0.7,
        1.4,
        5.35,
        4.65,
        [
            ("Рабочее определение", 18, INK, True),
            *bullets(
                [
                    "программный компонент внутри бизнес-процесса",
                    "получает входное событие и контекст обработки",
                    "возвращает результат в машиночитаемой форме",
                    "не ограничивается диалоговым ответом, как чат-бот",
                    "применяется для случаев, не полностью покрытых правилами",
                ],
                13.2,
            ),
        ],
    )
    table(
        s,
        6.35,
        1.55,
        [1.8, 2.2, 2.15],
        0.52,
        [
            ["Подход", "Основная функция", "Роль в работе"],
            ["Rule-based", "типовые случаи", "baseline"],
            ["Чат-бот", "диалог с пользователем", "не основной сценарий"],
            ["AI-агент", "решение внутри процесса", "основной объект"],
        ],
        font_size=10,
    )
    box(
        s,
        6.55,
        4.7,
        5.72,
        0.98,
        [("Ключевой критерий", 14, TEAL, True), ("агент участвует в исполнении процесса, а не только объясняет результат пользователю", 14.3, INK, False)],
        TEAL_LIGHT,
        TEAL,
    )


def slide5(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, 5)
    title(s, "Анализ подходов к AI-автоматизации")
    rows = [
        ("Rule-based", "скорость и контроль", "низкая адаптивность", BLUE_LIGHT, BLUE),
        ("ML", "покрытие сложных случаев", "нужны данные и обучение", GREEN_LIGHT, GREEN),
        ("LLM", "интерпретация контекста", "стоимость и латентность", AMBER_LIGHT, AMBER),
        ("Agent-based", "AI + backend-логика", "сложнее архитектурно", TEAL_LIGHT, TEAL),
    ]
    for i, (head, strong, weak, fill, line) in enumerate(rows):
        x = 0.72 + i * 3.05
        box(s, x, 1.45, 2.72, 2.25, [(head, 18, line, True), ("Сильная сторона", 10.5, MUTED, False), (strong, 13.5, INK, True), ("Ограничение", 10.5, MUTED, False), (weak, 12.5, INK, False)], fill, line)
    box(
        s,
        0.95,
        4.55,
        11.2,
        1.15,
        [("Вывод для прототипа", 18, INK, True), ("быстрый путь обработки сохраняется, а LLM используется только для ограниченного сегмента сложных случаев", 18, TEAL, True)],
        WHITE,
        TEAL,
        PP_ALIGN.CENTER,
    )


def slide6(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, 6)
    title(s, "Обоснование выбора технологического стека")
    components = [
        ("Spring Boot", "базовая сервисная архитектура", BLUE_LIGHT, BLUE),
        ("Spring AI", "интеграция с LLM через ChatClient", TEAL_LIGHT, TEAL),
        ("DeepSeek API", "внешний LLM-провайдер", PURPLE_LIGHT, PURPLE),
        ("Kafka", "асинхронная передача событий", AMBER_LIGHT, AMBER),
        ("Kotlin", "типобезопасность и JVM-совместимость", GREEN_LIGHT, GREEN),
        ("PostgreSQL", "хранение транзакций и результатов", WHITE, LINE),
    ]
    for idx, (name, purpose, fill, line) in enumerate(components):
        x = 0.75 + (idx % 3) * 4.05
        y = 1.45 + (idx // 3) * 1.65
        box(s, x, y, 3.45, 1.1, [(name, 18, line if line != LINE else INK, True), (purpose, 13.2, INK, False)], fill, line, PP_ALIGN.CENTER)
    box(
        s,
        0.95,
        5.65,
        11.1,
        0.72,
        [("Архитектурный принцип", 17, INK, True), ("прототип реализуется как enterprise-backend, а не как изолированный LLM-эксперимент", 15.5, MUTED, False)],
        WHITE,
        LINE,
        PP_ALIGN.CENTER,
    )


def slide7(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, 7)
    title(s, "Кейс исследования: классификация финансовых транзакций")
    box(
        s,
        0.7,
        1.35,
        4.55,
        3.65,
        [
            ("Почему этот кейс показателен", 18, INK, True),
            *bullets(
                [
                    "сочетает требования к скорости, точности и воспроизводимости",
                    "типовые транзакции проходят через rule-based слой",
                    "неоднозначные случаи направляются в LLM fallback",
                    "AI-компонент дополняет, а не заменяет процесс",
                ],
                13.2,
            ),
        ],
    )
    box(s, 5.85, 1.6, 2.1, 0.72, [("Сырая транзакция", 13, INK, True)], BLUE_LIGHT, BLUE, PP_ALIGN.CENTER)
    arrow(s, 8.05, 1.96, 8.55, 1.96)
    box(s, 8.7, 1.6, 2.55, 0.72, [("Rule-based", 13, INK, True)], WHITE, LINE, PP_ALIGN.CENTER)
    arrow(s, 9.95, 2.4, 9.95, 3.0)
    box(s, 8.55, 3.1, 2.8, 0.8, [("Confidence gate", 13, INK, True)], AMBER_LIGHT, AMBER, PP_ALIGN.CENTER)
    arrow(s, 8.45, 3.5, 7.55, 4.65)
    arrow(s, 11.45, 3.5, 11.9, 4.65)
    box(s, 6.25, 4.65, 2.7, 0.72, [("Финальная категория", 12.7, GREEN, True)], GREEN_LIGHT, GREEN, PP_ALIGN.CENTER)
    box(s, 10.4, 4.65, 2.0, 0.72, [("LLM fallback", 12.7, TEAL, True)], TEAL_LIGHT, TEAL, PP_ALIGN.CENTER)
    arrow(s, 10.4, 5.0, 9.05, 5.0)


def slide8(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, 8)
    title(s, "Архитектура прототипа")
    box(s, 0.7, 1.55, 1.65, 0.85, [("Client / API", 13, INK, True)], WHITE, LINE, PP_ALIGN.CENTER)
    arrow(s, 2.45, 1.98, 2.95, 1.98)
    box(s, 3.05, 1.55, 2.0, 0.85, [("Core Service", 13, BLUE, True), ("REST + orchestration", 9.5, MUTED, False)], BLUE_LIGHT, BLUE, PP_ALIGN.CENTER)
    arrow(s, 5.15, 1.98, 5.65, 1.98)
    box(s, 5.75, 1.55, 2.25, 0.85, [("Rule-based", 13, INK, True), ("fast path", 9.5, MUTED, False)], WHITE, LINE, PP_ALIGN.CENTER)
    arrow(s, 8.1, 1.98, 8.55, 1.98)
    box(s, 8.65, 1.55, 2.1, 0.85, [("Confidence", 13, AMBER, True), ("gate", 9.5, MUTED, False)], AMBER_LIGHT, AMBER, PP_ALIGN.CENTER)
    arrow(s, 9.7, 2.5, 9.7, 3.25)
    box(s, 8.45, 3.35, 2.55, 0.82, [("Kafka topic", 13, INK, True)], WHITE, LINE, PP_ALIGN.CENTER)
    arrow(s, 9.7, 4.25, 9.7, 5.0)
    box(s, 8.1, 5.1, 3.25, 0.82, [("LLM Classifier Agent", 13, TEAL, True)], TEAL_LIGHT, TEAL, PP_ALIGN.CENTER)
    box(s, 0.8, 4.85, 2.9, 0.9, [("PostgreSQL", 14, GREEN, True), ("transactions, labels, metadata", 10, MUTED, False)], GREEN_LIGHT, GREEN, PP_ALIGN.CENTER)
    arrow(s, 4.05, 2.4, 2.35, 4.75)
    arrow(s, 9.7, 5.95, 3.75, 5.3)
    box(
        s,
        0.75,
        6.35,
        11.6,
        0.42,
        [("LLM-компонент изолирован от быстрого пути и может масштабироваться отдельно", 13.5, INK, True)],
        WHITE,
        LINE,
        PP_ALIGN.CENTER,
    )


def slide9(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, 9)
    title(s, "Основной сценарий: Transaction Classifier Agent")
    steps = [
        ("1", "сырая транзакция"),
        ("2", "первичная категория + confidence"),
        ("3", "достаточная уверенность"),
        ("4", "низкая уверенность → LLM"),
        ("5", "уточнение и сохранение"),
    ]
    for i, (num, label) in enumerate(steps):
        x = 0.62 + i * 2.5
        box(s, x, 1.45, 2.0, 0.92, [(num, 19, TEAL, True), (label, 10.5, INK, True)], WHITE, LINE, PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            arrow(s, x + 2.05, 1.9, x + 2.38, 1.9)
    box(
        s,
        0.75,
        3.05,
        5.55,
        2.25,
        [
            ("Быстрый путь", 18, GREEN, True),
            *bullets(["результат принимается без обращения к LLM", "низкая задержка для большинства транзакций", "предсказуемая стоимость обработки"], 13.2),
        ],
        GREEN_LIGHT,
        GREEN,
    )
    box(
        s,
        6.75,
        3.05,
        5.55,
        2.25,
        [
            ("Расширенный путь", 18, TEAL, True),
            *bullets(["LLM используется только при low-confidence", "агент уточняет категорию по контексту", "результат возвращается в Core Service и БД"], 13.2),
        ],
        TEAL_LIGHT,
        TEAL,
    )
    pill(s, 1.45, 5.95, 2.35, "меньше UNDEFINED", GREEN_LIGHT, GREEN)
    pill(s, 4.15, 5.95, 2.35, "меньше ручной проверки", GREEN_LIGHT, GREEN)
    pill(s, 6.85, 5.95, 2.35, "быстрый основной поток", BLUE_LIGHT, BLUE)
    pill(s, 9.55, 5.95, 2.35, "контроль LLM-вызовов", AMBER_LIGHT, AMBER)


def slide10(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, 10)
    title(s, "Потенциал расширения: context-aware сценарий")
    box(
        s,
        0.75,
        1.45,
        4.7,
        3.9,
        [
            ("Перспектива развития", 18, INK, True),
            *bullets(
                [
                    "текущий агент использует только входные данные процесса",
                    "архитектура допускает получение дополнительного контекста",
                    "возможный tool: ограниченный доступ к истории транзакций",
                    "расширение не считается реализованной частью прототипа",
                ],
                13.4,
            ),
        ],
    )
    box(s, 6.15, 1.55, 2.1, 0.78, [("LLM Agent", 13.5, TEAL, True)], TEAL_LIGHT, TEAL, PP_ALIGN.CENTER)
    arrow(s, 8.35, 1.95, 9.25, 1.95)
    box(s, 9.4, 1.55, 2.35, 0.78, [("Достаточно контекста?", 12.2, INK, True)], AMBER_LIGHT, AMBER, PP_ALIGN.CENTER)
    arrow(s, 10.58, 2.42, 10.58, 3.25)
    box(s, 9.45, 3.35, 2.25, 0.78, [("Решение", 13.5, GREEN, True)], GREEN_LIGHT, GREEN, PP_ALIGN.CENTER)
    arrow(s, 9.3, 1.95, 7.2, 3.75)
    box(s, 6.05, 3.7, 2.3, 0.78, [("Запрос контекста", 12.5, BLUE, True)], BLUE_LIGHT, BLUE, PP_ALIGN.CENTER)
    arrow(s, 7.2, 4.55, 7.2, 5.15)
    box(s, 6.05, 5.25, 2.3, 0.78, [("База данных", 13.2, INK, True)], WHITE, LINE, PP_ALIGN.CENTER)
    arrow(s, 6.0, 5.65, 5.0, 2.0)
    pill(s, 9.05, 5.55, 2.8, "planned / not implemented", RED_LIGHT, RED)


def slide11(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, 11)
    title(s, "Методика экспериментальной оценки")
    box(s, 0.75, 1.35, 5.25, 1.2, [("Сравниваемые режимы", 17, INK, True), ("baseline: только rule-based классификация", 13.5, MUTED, False), ("hybrid: rule-based + LLM fallback для low-confidence", 13.5, TEAL, True)])
    table(
        s,
        6.45,
        1.35,
        [2.1, 4.0],
        0.58,
        [
            ["Группа", "Метрики"],
            ["Качество", "accuracy, UNDEFINED, manual review"],
            ["Производительность", "p95 latency, throughput, error rate"],
            ["Экономика", "cost per request, cost per flow"],
        ],
        font_size=10,
    )
    box(
        s,
        0.75,
        3.45,
        11.45,
        1.35,
        [("Данные эксперимента", 17, INK, True), ("синтетический сценарный набор транзакций, интерпретация результатов как демонстрационная оценка, а не production-метрики", 16, INK, False)],
        BLUE_LIGHT,
        BLUE,
        PP_ALIGN.CENTER,
    )
    pill(s, 1.2, 5.55, 2.3, "scenario dataset", WHITE, BLUE)
    arrow(s, 3.65, 5.73, 4.3, 5.73)
    pill(s, 4.45, 5.55, 2.3, "baseline run", WHITE, BLUE)
    arrow(s, 6.9, 5.73, 7.55, 5.73)
    pill(s, 7.7, 5.55, 2.3, "hybrid run", WHITE, TEAL)
    arrow(s, 10.15, 5.73, 10.8, 5.73)
    pill(s, 10.95, 5.55, 1.4, "compare", WHITE, GREEN)


def slide12(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, 12)
    title(s, "Предварительная демонстрационная оценка")
    metric_card(s, 0.75, 1.35, "Accuracy", "84-87%", "91-94%")
    metric_card(s, 0.75, 2.75, "UNDEFINED", "7-10%", "2-4%")
    metric_card(s, 0.75, 4.15, "Ручная проверка", "5-8%", "1-3%")
    box(
        s,
        7.25,
        1.35,
        4.65,
        1.25,
        [("Доля LLM-вызовов", 14, MUTED, False), ("5-12% потока", 26, AMBER, True), ("ограничение стоимости и задержки", 12, INK, False)],
        AMBER_LIGHT,
        AMBER,
        PP_ALIGN.CENTER,
    )
    box(
        s,
        7.25,
        3.15,
        4.65,
        1.85,
        [
            ("Интерпретация", 18, INK, True),
            *bullets(
                [
                    "rule-based слой остаётся быстрым",
                    "LLM закрывает low-confidence сегмент",
                    "значения являются модельной оценкой на синтетике",
                ],
                12.8,
            ),
        ],
    )
    box(s, 1.0, 6.08, 11.0, 0.48, [("Цифры не являются подтверждённым промышленным результатом и требуют экспериментальной проверки.", 13.5, RED, True)], RED_LIGHT, RED, PP_ALIGN.CENTER)


def slide13(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, 13)
    title(s, "Научная новизна, практическая значимость и выводы")
    box(
        s,
        0.68,
        1.35,
        3.75,
        4.5,
        [
            ("Научная новизна", 17, BLUE, True),
            *bullets(
                [
                    "AI-агент как управляемый компонент backend-процесса",
                    "гибридная обработка: rule-based + LLM fallback",
                    "интеграция LLM без перевода всего процесса на AI",
                ],
                12.7,
            ),
        ],
        BLUE_LIGHT,
        BLUE,
    )
    box(
        s,
        4.72,
        1.35,
        3.75,
        4.5,
        [
            ("Практическая значимость", 17, TEAL, True),
            *bullets(
                [
                    "воспроизводимый паттерн внедрения LLM",
                    "метрики качества, латентности и стоимости",
                    "event-driven изоляция долгих LLM-вызовов",
                ],
                12.7,
            ),
        ],
        TEAL_LIGHT,
        TEAL,
    )
    box(
        s,
        8.76,
        1.35,
        3.75,
        4.5,
        [
            ("Выводы", 17, GREEN, True),
            *bullets(
                [
                    "ценность даёт управляемое включение AI",
                    "hybrid сочетает скорость правил и гибкость LLM",
                    "прототип подтверждает инженерную реализуемость подхода",
                ],
                12.7,
            ),
        ],
        GREEN_LIGHT,
        GREEN,
    )


def slide14(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, 14)
    title(s, "Источники и технологическая база")
    left = [
        "работы по LLM-based agents и agentic AI",
        "исследования AI-автоматизации бизнес-процессов",
        "материалы по event-driven архитектурам",
    ]
    right = [
        "документация Spring AI",
        "документация Spring Boot и Spring Kafka",
        "документация DeepSeek API",
        "документация Apache Kafka",
        "документация Kotlin",
    ]
    box(s, 0.9, 1.55, 5.25, 3.8, [("Исследовательская база", 19, INK, True), *bullets(left, 15)], WHITE, LINE)
    box(s, 6.75, 1.55, 5.25, 3.8, [("Технологическая база", 19, INK, True), *bullets(right, 15)], BLUE_LIGHT, BLUE)
    box(
        s,
        1.1,
        6.05,
        10.7,
        0.55,
        [("Спасибо за внимание", 20, TEAL, True)],
        WHITE,
        TEAL,
        PP_ALIGN.CENTER,
    )


NOTES = [
    ("Разработка и внедрение AI-агента для оптимизации бизнес-процессов", "Представить тему как исследование управляемого внедрения AI-агента в backend-процесс. Подчеркнуть, что основной фокус не на чат-интерфейсе, а на измеримом изменении процесса."),
    ("Актуальность и проблема исследования", "Объяснить конфликт между rule-based скоростью и LLM-гибкостью. Основная проблема: full LLM слишком дорогой и медленный, но правила не закрывают неоднозначные случаи."),
    ("Объект, предмет, цель, задачи и методы", "Коротко пройти по формальным элементам исследования. Цель связать с гибридным подходом rule-based + LLM fallback."),
    ("AI-агент в контексте работы", "Зафиксировать определение агента: компонент процесса, который получает событие, использует контекст и возвращает машиночитаемое решение."),
    ("Анализ подходов к AI-автоматизации", "Показать, почему выбран agent-based hybrid: он сохраняет быстрый путь и добавляет AI только там, где правила не дают уверенного результата."),
    ("Обоснование выбора технологического стека", "Пояснить, что Spring/Kotlin выбраны из-за enterprise-backend контекста, Kafka изолирует долгие LLM-вызовы, PostgreSQL хранит результаты и атрибуты обработки."),
    ("Кейс исследования: классификация финансовых транзакций", "Показать сценарий как удобный исследовательский кейс: есть скорость, точность, воспроизводимость и понятный low-confidence сегмент."),
    ("Архитектура прототипа", "Объяснить поток: Core Service принимает событие, rule-based слой пытается классифицировать, low-confidence случаи уходят через Kafka к LLM-агенту, результаты сохраняются в PostgreSQL."),
    ("Основной сценарий: Transaction Classifier Agent", "Подчеркнуть, что это главный агент работы. Он снижает долю неопределённых транзакций и ручной проверки, но не замедляет основной поток."),
    ("Потенциал расширения: context-aware сценарий", "Отдельно сказать, что получение контекста из БД является перспективой развития, а не реализованной частью текущего прототипа."),
    ("Методика экспериментальной оценки", "Раскрыть baseline и hybrid режимы. Все результаты нужно читать как демонстрационные, потому что набор транзакций синтетический."),
    ("Предварительная демонстрационная оценка", "Пояснить, что диапазоны являются ориентирами для эксперимента. Их нельзя выдавать за production-метрики."),
    ("Научная новизна, практическая значимость и выводы", "Сформулировать итог: управляемый гибридный подход предпочтительнее полного перевода процесса на LLM."),
    ("Источники и технологическая база", "Закрыть выступление ссылкой на исследовательскую и технологическую базу: LLM agents, Spring AI, Kafka, Kotlin, PostgreSQL, DeepSeek API."),
]


def write_notes() -> None:
    lines = [
        "# FinSense AI agent description: speaker notes",
        "",
        "Заметки подготовлены по `docs/presentation/DESCRIPTION.md`. Демонстрационные метрики на слайде 12 следует явно отделять от production-результатов.",
        "",
    ]
    for idx, (heading, note) in enumerate(NOTES, start=1):
        lines.extend([f"## Слайд {idx}. {heading}", "", note, ""])
    NOTES_PATH.write_text("\n".join(lines), encoding="utf-8")


def build() -> None:
    prs = Presentation()
    prs.slide_width = Inches(WIDE_W)
    prs.slide_height = Inches(WIDE_H)
    for maker in (
        slide1,
        slide2,
        slide3,
        slide4,
        slide5,
        slide6,
        slide7,
        slide8,
        slide9,
        slide10,
        slide11,
        slide12,
        slide13,
        slide14,
    ):
        maker(prs)
    prs.save(PPTX_PATH)
    write_notes()
    print(PPTX_PATH)
    print(NOTES_PATH)


if __name__ == "__main__":
    build()
