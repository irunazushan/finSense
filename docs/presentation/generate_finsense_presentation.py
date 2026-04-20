from __future__ import annotations

from pathlib import Path
import sys

LOCAL_DEPS = Path(__file__).resolve().parent / "_pptx_deps"
if LOCAL_DEPS.exists():
    sys.path.insert(0, str(LOCAL_DEPS))

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt
except ModuleNotFoundError as exc:
    raise SystemExit(
        "Install python-pptx first: python -m pip install python-pptx"
    ) from exc


OUT_DIR = Path(__file__).resolve().parent
PPTX_PATH = OUT_DIR / "FinSense_AI_business_process.pptx"

BG = RGBColor(7, 26, 45)
PANEL = RGBColor(13, 42, 66)
PANEL_2 = RGBColor(18, 59, 88)
TEXT = RGBColor(245, 247, 250)
MUTED = RGBColor(167, 190, 211)
ACCENT = RGBColor(45, 212, 191)
ACCENT_2 = RGBColor(245, 158, 11)
GREEN = RGBColor(34, 197, 94)
RED = RGBColor(251, 113, 133)

FONT = "Arial"


def add_background(slide):
    bg = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()

    stripe = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(0.18), Inches(7.5)
    )
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = ACCENT
    stripe.line.fill.background()

    depth = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(10.5), 0, Inches(2.83), Inches(7.5)
    )
    depth.fill.solid()
    depth.fill.fore_color.rgb = RGBColor(9, 34, 57)
    depth.line.fill.background()


def text_box(slide, x, y, w, h, paragraphs, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for index, item in enumerate(paragraphs):
        text, size, color, bold = item
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.text = text
        p.alignment = align
        p.space_after = Pt(4)
        run = p.runs[0]
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
    return box


def card(slide, x, y, w, h, paragraphs, fill=PANEL, line=RGBColor(31, 86, 116), align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1.2)
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.08)
    tf.word_wrap = True
    for index, item in enumerate(paragraphs):
        text, size, color, bold = item
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.text = text
        p.alignment = align
        p.space_after = Pt(4)
        run = p.runs[0]
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
    return shape


def title(slide, text, subtitle=None):
    paragraphs = [(text, 32, TEXT, True)]
    if subtitle:
        paragraphs.append((subtitle, 15, MUTED, False))
    text_box(slide, 0.55, 0.42, 10.8, 1.05, paragraphs)


def bullets(lines, size=18):
    return [(f"• {line}", size, TEXT, False) for line in lines]


def arrow(slide, x1, y1, x2, y2, color=ACCENT):
    shape = slide.shapes.add_connector(
        1, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    shape.line.color.rgb = color
    shape.line.width = Pt(2.3)
    return shape


def add_slides(prs: Presentation):
    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)
    add_background(slide)
    text_box(
        slide,
        0.65,
        0.95,
        8.8,
        1.35,
        [
            ("FinSense", 54, TEXT, True),
            ("LLM как инструмент оптимизации бизнес-процесса", 25, ACCENT, False),
        ],
    )
    card(slide, 0.72, 2.85, 4.0, 1.05, [("10 минут", 16, MUTED, False), ("магистерская конференция", 22, TEXT, True)], PANEL_2)
    card(slide, 5.05, 2.85, 4.1, 1.05, [("Ключевая идея", 16, MUTED, False), ("LLM усиливает сложные участки процесса", 20, TEXT, True)])
    card(slide, 9.45, 2.85, 3.0, 1.05, [("Формат", 16, MUTED, False), ("синтетический бизнес-кейс", 20, TEXT, True)], RGBColor(23, 51, 74))
    text_box(slide, 0.72, 5.8, 11.4, 0.6, [("Демонстрация: классификация транзакций, персональные рекомендации, event-driven архитектура", 18, MUTED, False)])

    slide = prs.slides.add_slide(blank)
    add_background(slide)
    title(slide, "Проблема бизнеса")
    card(slide, 0.72, 1.65, 3.8, 3.25, [("Категоризация", 24, ACCENT, True), *bullets(["правила и MCC не всегда отражают смысл платежа", "новые мерчанты и смешанные описания дают ошибки", "часть операций уходит в UNDEFINED"], 17)])
    card(slide, 4.8, 1.65, 3.8, 3.25, [("Рекомендации", 24, ACCENT_2, True), *bullets(["часто шаблонные", "не используют реальное поведение клиента", "трудно объяснить, на каких данных основан совет"], 17)])
    card(slide, 8.88, 1.65, 3.8, 3.25, [("Операции", 24, RED, True), *bullets(["ручная проверка стоит дорого", "LLM для всех транзакций медленная и дорогая", "нужен баланс качества, latency и стоимости"], 17)])
    text_box(slide, 0.78, 5.45, 11.5, 0.8, [("Задача: повысить качество решений, не превращая каждую операцию в дорогой LLM-вызов.", 24, TEXT, True)], PP_ALIGN.CENTER)

    slide = prs.slides.add_slide(blank)
    add_background(slide)
    title(slide, "Идея решения: гибридный AI-пайплайн")
    card(slide, 0.7, 1.7, 3.2, 1.35, [("1. Быстрый слой", 18, MUTED, False), ("Rule/ML-first классификация", 23, TEXT, True), ("< 50 мс для типовых операций", 15, ACCENT, False)])
    arrow(slide, 4.05, 2.37, 4.85, 2.37)
    card(slide, 4.95, 1.7, 3.25, 1.35, [("2. Confidence gate", 18, MUTED, False), ("Порог уверенности 0.9", 23, TEXT, True), ("только сложные случаи идут дальше", 15, ACCENT, False)], PANEL_2)
    arrow(slide, 8.35, 2.37, 9.15, 2.37)
    card(slide, 9.25, 1.7, 3.25, 1.35, [("3. LLM fallback", 18, MUTED, False), ("Контекстная классификация", 23, TEXT, True), ("история + описание + merchant", 15, ACCENT, False)])
    card(slide, 1.15, 4.35, 5.25, 1.35, [("Financial Coach Agent", 24, TEXT, True), ("LLM формирует рекомендации на основе инструментов анализа расходов", 18, MUTED, False)], RGBColor(16, 47, 70))
    card(slide, 6.9, 4.35, 5.25, 1.35, [("Kafka + Postgres", 24, TEXT, True), ("тяжёлые операции асинхронны, результат сохраняется и доступен через API", 18, MUTED, False)], RGBColor(16, 47, 70))

    slide = prs.slides.add_slide(blank)
    add_background(slide)
    title(slide, "Архитектура FinSense")
    card(slide, 0.7, 1.35, 2.55, 1.1, [("Core Service", 22, TEXT, True), ("оркестратор, REST API, доменная модель", 14, MUTED, False)], PANEL_2)
    card(slide, 3.65, 1.35, 2.55, 1.1, [("Classifier", 22, TEXT, True), ("быстрая rule/ML классификация", 14, MUTED, False)])
    card(slide, 6.6, 1.35, 2.55, 1.1, [("LLM Classifier", 22, TEXT, True), ("fallback для неуверенных транзакций", 14, MUTED, False)])
    card(slide, 9.55, 1.35, 2.55, 1.1, [("Coach Agent", 22, TEXT, True), ("персональные советы по расходам", 14, MUTED, False)])
    card(slide, 1.55, 3.45, 4.6, 1.35, [("Kafka", 30, ACCENT, True), ("raw-transactions, llm-classifier-requests/responses, coach-requests/responses", 16, MUTED, False)], RGBColor(8, 37, 59), align=PP_ALIGN.CENTER)
    card(slide, 7.15, 3.45, 4.6, 1.35, [("PostgreSQL", 30, ACCENT_2, True), ("транзакции, статусы, рекомендации, JSONB с AI-метаданными", 16, MUTED, False)], RGBColor(8, 37, 59), align=PP_ALIGN.CENTER)
    text_box(slide, 0.8, 5.8, 11.6, 0.65, [("LLM вынесена в отдельные event-driven компоненты и может масштабироваться независимо.", 21, TEXT, True)], PP_ALIGN.CENTER)

    slide = prs.slides.add_slide(blank)
    add_background(slide)
    title(slide, "Поток классификации транзакции")
    for idx, (x, label) in enumerate([(0.65, "raw transaction"), (3.7, "Core сохраняет"), (6.75, "Classifier"), (9.8, "confidence")], start=1):
        card(slide, x, 1.65, 2.2, 1.0, [(str(idx), 26, ACCENT, True), (label, 17, TEXT, True)], align=PP_ALIGN.CENTER)
    arrow(slide, 2.95, 2.15, 3.6, 2.15)
    arrow(slide, 6.0, 2.15, 6.65, 2.15)
    arrow(slide, 9.05, 2.15, 9.7, 2.15)
    card(slide, 1.15, 4.0, 5.1, 1.25, [("Если confidence >= 0.9", 25, GREEN, True), ("транзакция сразу получает статус CLASSIFIED, источник ML/RULE", 18, TEXT, False)], RGBColor(14, 53, 46), GREEN)
    card(slide, 7.05, 4.0, 5.1, 1.25, [("Если confidence < 0.9", 25, ACCENT_2, True), ("Core публикует событие в llm-classifier-requests и ждёт ответ агента", 18, TEXT, False)], RGBColor(58, 43, 13), ACCENT_2)
    text_box(slide, 0.85, 6.05, 11.4, 0.5, [("Состояния в БД: ML_CLASSIFYING → LLM_CLASSIFYING → CLASSIFIED / FAILED", 19, MUTED, False)], PP_ALIGN.CENTER)

    slide = prs.slides.add_slide(blank)
    add_background(slide)
    title(slide, "Как LLM улучшает процесс")
    card(slide, 0.75, 1.55, 3.7, 3.75, [("До", 32, RED, True), *bullets(["неоднозначные описания попадают в UNDEFINED", "ошибки MCC и keyword-правил остаются финальными", "оператор вручную разбирает спорные случаи"], 18)], RGBColor(51, 24, 39), RED)
    card(slide, 4.85, 1.55, 3.7, 3.75, [("После", 32, GREEN, True), *bullets(["LLM видит текущую операцию и историю", "возвращает строго допустимую категорию", "даёт confidence и короткое reasoning"], 18)], RGBColor(17, 53, 36), GREEN)
    card(slide, 8.95, 1.55, 3.7, 3.75, [("Оптимизация", 32, ACCENT, True), *bullets(["LLM вызывается только для 5-10% сложных операций", "основной поток остаётся быстрым", "затраты контролируются threshold-порогом"], 18)], PANEL, ACCENT)
    text_box(slide, 0.9, 5.95, 11.25, 0.55, [("Главная ценность: LLM закрывает «серую зону» правил, но не становится узким местом всей системы.", 21, TEXT, True)], PP_ALIGN.CENTER)

    slide = prs.slides.add_slide(blank)
    add_background(slide)
    title(slide, "Financial Coach Agent: рекомендации на фактах")
    cards = [
        ("Tool 1", "Spending by category", "структура расходов"),
        ("Tool 2", "Monthly delta", "рост/снижение по периодам"),
        ("Tool 3", "Top merchants", "главные источники трат"),
        ("Tool 4", "Spikes", "аномальные всплески"),
    ]
    for x, (cap, main, sub) in zip([0.65, 3.75, 6.85, 9.95], cards):
        card(slide, x, 1.5, 2.75, 1.35, [(cap, 17, MUTED, False), (main, 22, TEXT, True), (sub, 15, ACCENT, False)])
    card(slide, 1.05, 4.1, 5.1, 1.45, [("Детерминированная часть", 24, ACCENT_2, True), ("SQL-агрегации дают LLM проверяемые факты: суммы, категории, мерчантов, даты.", 18, TEXT, False)], RGBColor(58, 43, 13), ACCENT_2)
    card(slide, 7.05, 4.1, 5.1, 1.45, [("Генеративная часть", 24, ACCENT, True), ("LLM превращает факты в понятный пользователю совет на русском языке.", 18, TEXT, False)], RGBColor(14, 53, 51), ACCENT)

    slide = prs.slides.add_slide(blank)
    add_background(slide)
    title(slide, "Метрики эффекта", "пилотная оценка MVP на синтетических данных")
    metric_cards = [
        (0.75, 1.55, 3.05, "Точность", "85-90% → 95%+", "rule-only → hybrid", GREEN),
        (4.05, 1.55, 3.05, "LLM fallback", "5-10%", "только сложные случаи", ACCENT),
        (7.35, 1.55, 2.35, "Manual review", ">5% → <1%", "меньше UNDEFINED", GREEN),
        (9.95, 1.55, 2.35, "Latency p95", "<2 c", "с учётом LLM", ACCENT_2),
        (0.75, 3.35, 3.05, "Стоимость", "+10-15%", "не +100%, потому что LLM не для всех", ACCENT_2),
        (4.05, 3.35, 3.05, "Рекомендации", "80%+", "целевой positive feedback", GREEN),
    ]
    for x, y, w, top, value, sub, color in metric_cards:
        card(slide, x, y, w, 1.1, [(top, 16, MUTED, False), (value, 29, color, True), (sub, 13, TEXT, False)])
    card(slide, 7.35, 3.35, 4.95, 1.1, [("Наблюдаемость", 16, MUTED, False), ("requestId, transactionId, tokens, latency, model", 21, TEXT, True), ("аудит AI-решений и воспроизводимость анализа", 14, ACCENT, False)])
    text_box(slide, 0.85, 5.75, 11.35, 0.8, [("Метрики нужно валидировать на размеченном наборе: базовый rule-only прогон, затем гибридный прогон с LLM fallback.", 20, MUTED, False)], PP_ALIGN.CENTER)

    slide = prs.slides.add_slide(blank)
    add_background(slide)
    title(slide, "Итог: синтетический кейс, переносимый подход")
    card(slide, 0.75, 1.35, 5.55, 1.5, [("Важно", 24, ACCENT_2, True), ("FinSense — синтетический бизнес-кейс, реализованный специально для демонстрации внедрения ИИ в backend-бизнес-процесс.", 19, TEXT, False)], RGBColor(58, 43, 13), ACCENT_2)
    card(slide, 6.85, 1.35, 5.55, 1.5, [("Обобщение", 24, ACCENT, True), ("Паттерн переносим: быстрый детерминированный слой + LLM для сложных случаев + аудит + асинхронная обработка.", 19, TEXT, False)], RGBColor(14, 53, 51), ACCENT)
    examples = [
        ("Страхование", "первичная обработка заявлений и документов по страховым случаям"),
        ("Кредитование", "предварительный анализ кредитной заявки и объяснение отказов"),
        ("Support", "маршрутизация обращений, приоритизация и подготовка ответа"),
        ("Документы", "анализ договоров, счетов, актов и первички"),
    ]
    for x, (head, body) in zip([0.75, 3.95, 7.15, 10.35], examples):
        card(slide, x, 3.35, 2.8 if x < 10 else 2.2, 1.55, [(head, 22, TEXT, True), (body, 15, MUTED, False)])
    text_box(slide, 1.0, 5.95, 11.0, 0.55, [("Вывод: LLM даёт максимальный эффект там, где есть неоднозначность, контекст и высокая стоимость ручного решения.", 22, TEXT, True)], PP_ALIGN.CENTER)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    add_slides(prs)
    prs.save(PPTX_PATH)
    print(PPTX_PATH)


if __name__ == "__main__":
    main()
