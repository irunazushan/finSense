from __future__ import annotations

from pathlib import Path
import sys

LOCAL_DEPS = Path(__file__).resolve().parent / "_pptx_deps"
if LOCAL_DEPS.exists():
    sys.path.insert(0, str(LOCAL_DEPS))

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT_DIR = Path(__file__).resolve().parent
PPTX_PATH = OUT_DIR / "FinSense_thesis_predefense.pptx"

NAVY = RGBColor(12, 39, 74)
NAVY_2 = RGBColor(22, 63, 112)
BLUE = RGBColor(37, 99, 235)
LIGHT_BLUE = RGBColor(232, 241, 255)
ICE = RGBColor(247, 250, 255)
WHITE = RGBColor(255, 255, 255)
MUTED = RGBColor(84, 103, 128)
GRID = RGBColor(204, 219, 239)
GREEN = RGBColor(21, 128, 61)
ORANGE = RGBColor(194, 101, 10)
RED = RGBColor(185, 28, 28)
FONT = "Arial"


def add_bg(slide):
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = ICE
    bg.line.fill.background()

    top = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.15))
    top.fill.solid()
    top.fill.fore_color.rgb = NAVY
    top.line.fill.background()

    side = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(12.85), 0, Inches(0.48), Inches(7.5))
    side.fill.solid()
    side.fill.fore_color.rgb = LIGHT_BLUE
    side.line.fill.background()

    mark = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(10.9), Inches(-0.85), Inches(2.8), Inches(2.8))
    mark.fill.solid()
    mark.fill.fore_color.rgb = RGBColor(221, 235, 255)
    mark.line.fill.background()


def set_text(run, size, color=NAVY, bold=False):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold


def box(slide, x, y, w, h, items, align=PP_ALIGN.LEFT, fill=None, line=None, rounded=False):
    if fill is None:
        shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    else:
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if rounded else MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(x),
            Inches(y),
            Inches(w),
            Inches(h),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        shape.line.color.rgb = line or GRID
        shape.line.width = Pt(1.1)

    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.14)
    tf.margin_right = Inches(0.14)
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.08)

    for i, (text, size, color, bold) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.alignment = align
        p.space_after = Pt(3)
        set_text(p.runs[0], size, color, bold)
    return shape


def title(slide, text, subtitle=None):
    box(slide, 0.55, 0.42, 11.6, 0.65, [(text, 29, NAVY, True)])
    if subtitle:
        box(slide, 0.58, 1.05, 10.9, 0.36, [(subtitle, 13, MUTED, False)])


def bullets(lines, size=16, color=NAVY):
    return [(f"• {line}", size, color, False) for line in lines]


def pill(slide, x, y, w, text, fill=LIGHT_BLUE, color=NAVY):
    return box(slide, x, y, w, 0.45, [(text, 13, color, True)], PP_ALIGN.CENTER, fill, GRID, True)


def arrow(slide, x1, y1, x2, y2, color=BLUE):
    conn = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(2.0)
    return conn


def metric_bar(slide, x, y, label, before, after, unit="%", max_value=100):
    box(slide, x, y, 2.6, 0.35, [(label, 13, NAVY, True)])
    box(slide, x, y + 0.48, 1.0, 0.32, [("до", 11, MUTED, False)], PP_ALIGN.RIGHT)
    before_w = 3.8 * before / max_value
    after_w = 3.8 * after / max_value
    b = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x + 1.15), Inches(y + 0.5), Inches(before_w), Inches(0.24))
    b.fill.solid()
    b.fill.fore_color.rgb = RGBColor(163, 180, 204)
    b.line.fill.background()
    box(slide, x + 5.1, y + 0.43, 0.85, 0.32, [(f"{before:g}{unit}", 11, MUTED, True)])

    box(slide, x, y + 0.88, 1.0, 0.32, [("после", 11, NAVY, False)], PP_ALIGN.RIGHT)
    a = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x + 1.15), Inches(y + 0.9), Inches(after_w), Inches(0.24))
    a.fill.solid()
    a.fill.fore_color.rgb = BLUE if after >= before else GREEN
    a.line.fill.background()
    box(slide, x + 5.1, y + 0.83, 0.85, 0.32, [(f"{after:g}{unit}", 11, NAVY, True)])


def build_slides(prs: Presentation):
    blank = prs.slide_layouts[6]

    # 1
    s = prs.slides.add_slide(blank)
    add_bg(s)
    box(s, 0.7, 0.88, 10.6, 1.45, [("Разработка и внедрение AI-агента", 34, NAVY, True), ("для оптимизации бизнес-процессов", 30, NAVY_2, True)])
    box(s, 0.74, 2.55, 9.7, 0.75, [("Анализ существующих решений и разработка прототипа на базе Spring и Kotlin", 21, MUTED, False)])
    box(s, 0.72, 4.12, 3.5, 1.05, [("Формат", 14, MUTED, False), ("предзащита магистерской диссертации", 21, NAVY, True)], fill=WHITE, rounded=True)
    box(s, 4.55, 4.12, 3.25, 1.05, [("Прототип", 14, MUTED, False), ("FinSense", 25, BLUE, True)], fill=WHITE, rounded=True)
    box(s, 8.15, 4.12, 3.6, 1.05, [("Фокус", 14, MUTED, False), ("оптимизация transaction flow", 20, NAVY, True)], fill=WHITE, rounded=True)
    box(s, 0.74, 6.25, 11.4, 0.5, [("Основной агент: Transaction Classifier Agent. Дополнительный сценарий: Financial Coach Agent.", 18, NAVY, True)], PP_ALIGN.CENTER)

    # 2
    s = prs.slides.add_slide(blank)
    add_bg(s)
    title(s, "Введение: контекст работы")
    box(s, 0.75, 1.45, 3.6, 3.8, [("Что меняется", 22, NAVY, True), *bullets(["LLM стали применимы не только для текста, но и для задач принятия решений", "tool calling позволяет подключать модель к API, БД и бизнес-данным", "AI-компоненты можно встраивать в существующие backend-процессы"], 16)], fill=WHITE, rounded=True)
    box(s, 4.85, 1.45, 3.6, 3.8, [("Что важно бизнесу", 22, NAVY, True), *bullets(["меньше ручного труда", "быстрее обработка операций", "выше качество решений", "измеримый эффект: latency, cost, accuracy"], 16)], fill=WHITE, rounded=True)
    box(s, 8.95, 1.45, 3.6, 3.8, [("Что исследуется", 22, NAVY, True), *bullets(["как встроить AI-агента в процесс", "как ограничить LLM контрактами", "как сохранить контролируемость и поддержку в JVM-стеке"], 16)], fill=WHITE, rounded=True)
    box(s, 0.95, 5.95, 11.0, 0.55, [("Таким образом, диплом рассматривает AI-агента как инженерный компонент бизнес-процесса, а не как отдельный чат-интерфейс.", 18, NAVY, True)], PP_ALIGN.CENTER)

    # 3
    s = prs.slides.add_slide(blank)
    add_bg(s)
    title(s, "Проблема, цель и объект исследования")
    box(s, 0.75, 1.45, 3.65, 3.85, [("Проблема", 23, NAVY, True), *bullets(["типовые правила быстрые, но ошибаются в неоднозначных случаях", "LLM для всех операций дорогой и медленный", "нужен управляемый fallback, а не замена всего процесса"], 16)], fill=WHITE, rounded=True)
    box(s, 4.85, 1.45, 3.65, 3.85, [("Цель", 23, NAVY, True), ("Разработать прототип AI-агента, который оптимизирует процесс классификации финансовых транзакций.", 18, NAVY, False), ("Дополнительно показать расширение архитектуры через агента финансовых рекомендаций.", 17, MUTED, False)], fill=WHITE, rounded=True)
    box(s, 9.0, 1.45, 3.35, 3.85, [("Предмет", 23, NAVY, True), *bullets(["LLM-agent внутри backend-сервиса", "Spring/Kotlin интеграция", "Kafka-события", "метрики качества и стоимости"], 16)], fill=WHITE, rounded=True)
    pill(s, 1.1, 5.9, 2.45, "transaction flow")
    pill(s, 3.85, 5.9, 2.45, "LLM fallback")
    pill(s, 6.6, 5.9, 2.45, "AI agent")
    pill(s, 9.35, 5.9, 2.45, "measurable effect")

    # 4
    s = prs.slides.add_slide(blank)
    add_bg(s)
    title(s, "Глава 1. Теоретические основы")
    chapters = [
        ("1.1", "Бизнес-процесс и оптимизация", "что оптимизируется: время, стоимость, качество, доля ручной работы"),
        ("1.2", "Эволюция автоматизации", "правила, RPA, ML и LLM как разные уровни автоматизации"),
        ("1.3", "AI-агенты и LLM", "не все AI-агенты обязаны быть LLM-based, но в этой работе LLM является ядром агента"),
        ("1.4", "Метрики эффективности", "accuracy, fallback rate, latency, cost, usefulness"),
        ("1.5", "Java/Kotlin и Spring", "почему AI удобно встраивать в существующий JVM backend"),
    ]
    for y, (num, head, body) in zip([1.35, 2.25, 3.15, 4.05, 4.95], chapters):
        box(s, 0.85, y, 1.0, 0.68, [(num, 18, WHITE, True)], PP_ALIGN.CENTER, NAVY, NAVY, True)
        box(s, 2.05, y, 9.85, 0.68, [(head, 18, NAVY, True), (body, 13, MUTED, False)], fill=WHITE, rounded=True)
    box(s, 0.95, 6.12, 11.0, 0.55, [("На выступлении этот слайд задаёт карту первой главы: от общих бизнес-процессов к конкретной технологии внедрения.", 17, NAVY, True)], PP_ALIGN.CENTER)

    # 5
    s = prs.slides.add_slide(blank)
    add_bg(s)
    title(s, "Что в работе называется AI-агентом")
    box(s, 0.8, 1.55, 2.35, 1.1, [("Событие", 22, NAVY, True), ("запрос из бизнес-процесса", 14, MUTED, False)], PP_ALIGN.CENTER, WHITE, GRID, True)
    arrow(s, 3.2, 2.1, 3.85, 2.1)
    box(s, 4.0, 1.55, 2.35, 1.1, [("Контекст", 22, NAVY, True), ("данные транзакции или пользователя", 14, MUTED, False)], PP_ALIGN.CENTER, WHITE, GRID, True)
    arrow(s, 6.4, 2.1, 7.05, 2.1)
    box(s, 7.2, 1.55, 2.35, 1.1, [("LLM + правила", 22, NAVY, True), ("ограниченная предметная область", 14, MUTED, False)], PP_ALIGN.CENTER, WHITE, GRID, True)
    arrow(s, 9.6, 2.1, 10.25, 2.1)
    box(s, 10.4, 1.55, 2.05, 1.1, [("Решение", 22, BLUE, True), ("категория или совет", 14, MUTED, False)], PP_ALIGN.CENTER, WHITE, GRID, True)
    box(s, 0.9, 3.35, 5.35, 1.75, [("Transaction Classifier Agent", 22, NAVY, True), *bullets(["специализированный AI-агент", "получает low-confidence транзакцию", "возвращает категорию и confidence", "оптимизирует основной transaction flow"], 15)], fill=LIGHT_BLUE, rounded=True)
    box(s, 6.75, 3.35, 5.35, 1.75, [("Financial Coach Agent", 22, NAVY, True), *bullets(["дополнительный агентный сценарий", "использует tools анализа расходов", "формирует персональную рекомендацию", "показывает расширяемость архитектуры"], 15)], fill=WHITE, rounded=True)
    box(s, 1.0, 6.0, 10.9, 0.55, [("В этой работе агент не обязан быть полностью автономным: достаточно, что он принимает контекстное решение и возвращает его в бизнес-процесс.", 17, NAVY, True)], PP_ALIGN.CENTER)

    # 6
    s = prs.slides.add_slide(blank)
    add_bg(s)
    title(s, "Почему Java/Kotlin и Spring важны для темы")
    box(s, 0.75, 1.45, 3.6, 3.8, [("Enterprise-контекст", 22, NAVY, True), *bullets(["Java/JVM часто уже присутствуют в backend-ландшафте", "AI проще внедрять рядом с существующими сервисами", "меньше межъязыковых интеграций и runtime-слоёв"], 16)], fill=WHITE, rounded=True)
    box(s, 4.85, 1.45, 3.6, 3.8, [("Kotlin", 22, NAVY, True), *bullets(["совместим с Java-библиотеками", "меньше шаблонного кода", "удобен для DTO, сервисов и конфигурации"], 16)], fill=WHITE, rounded=True)
    box(s, 8.95, 1.45, 3.6, 3.8, [("Spring", 22, NAVY, True), *bullets(["Spring Boot для микросервисов", "Spring Kafka для событий", "Spring AI для ChatClient и tools", "Actuator/Micrometer для метрик"], 16)], fill=WHITE, rounded=True)
    box(s, 0.9, 5.95, 11.2, 0.55, [("Обоснование выбора: стек понятен Java/Kotlin-командам и допускает тонкую настройку интеграций, логирования, retry и контрактов.", 18, NAVY, True)], PP_ALIGN.CENTER)

    # 7
    s = prs.slides.add_slide(blank)
    add_bg(s)
    title(s, "Глава 2. Анализ существующих AI-внедрений")
    cases = [
        ("Клиентская поддержка", "+14% productivity в NBER-исследовании; Klarna: 2/3 чатов AI assistant"),
        ("Страхование / claims", "автоматизация извлечения данных, triage и первичной оценки"),
        ("Кредитование", "ускорение предварительной проверки и анализа документов"),
        ("Документооборот", "извлечение реквизитов, классификация, поиск несоответствий"),
        ("Разработка ПО", "GitHub Copilot: 55% faster completion в эксперименте"),
    ]
    for y, (head, body) in zip([1.35, 2.25, 3.15, 4.05, 4.95], cases):
        box(s, 0.85, y, 3.1, 0.68, [(head, 17, WHITE, True)], PP_ALIGN.CENTER, NAVY, NAVY, True)
        box(s, 4.2, y, 7.7, 0.68, [(body, 16, NAVY, False)], fill=WHITE, rounded=True)
    box(s, 0.9, 6.05, 11.1, 0.58, [("Задача главы: показать измеримые эффекты AI в реальных процессах и вывести требования к собственному прототипу.", 18, NAVY, True)], PP_ALIGN.CENTER)

    # 8
    s = prs.slides.add_slide(blank)
    add_bg(s)
    title(s, "Сравнение технологий реализации")
    headers = ["Подход", "Сильная сторона", "Ограничение", "Роль в работе"]
    widths = [2.15, 3.2, 3.2, 3.0]
    xs = [0.6, 2.85, 6.18, 9.5]
    for x, w, h in zip(xs, widths, headers):
        box(s, x, 1.38, w, 0.5, [(h, 13, WHITE, True)], PP_ALIGN.CENTER, NAVY, NAVY)
    rows = [
        ("API языковых моделей", "быстрый доступ к сильным моделям", "стоимость, задержка, зависимость от провайдера", "внешний слой вывода"),
        ("Python-фреймворки агентов", "богатая экосистема экспериментов", "отдельный стек и runtime", "исследования и быстрые прототипы"),
        ("Spring AI", "нативно для Spring-приложений", "экосистема моложе Python-стека", "основа прототипа"),
        ("Apache Kafka", "асинхронность и масштабирование", "нужны схемы и идемпотентность", "шина событий"),
        ("PostgreSQL JSONB", "структурные данные + гибкость AI-метаданных", "нужен контроль формата", "хранилище результатов"),
    ]
    y = 1.95
    for row in rows:
        for x, w, text in zip(xs, widths, row):
            box(s, x, y, w, 0.72, [(text, 12, NAVY, False)], PP_ALIGN.CENTER, WHITE, GRID)
        y += 0.78
    box(s, 0.85, 6.2, 11.4, 0.52, [("Вывод: для JVM-микросервисов рационально внедрять AI-агента в Spring/Kotlin-стек, а долгие LLM-вызовы выносить в event-driven обработку.", 17, NAVY, True)], PP_ALIGN.CENTER)

    # 9
    s = prs.slides.add_slide(blank)
    add_bg(s)
    title(s, "Глава 3. Прототип FinSense")
    box(s, 0.8, 1.45, 3.5, 3.8, [("Синтетический кейс", 23, NAVY, True), *bullets(["не production banking system", "создан для демонстрации AI-внедрения", "моделирует обработку транзакций и рекомендации"], 16)], fill=WHITE, rounded=True)
    box(s, 4.85, 1.45, 3.5, 3.8, [("Основная оптимизация", 23, NAVY, True), *bullets(["transaction classification flow", "low-confidence транзакции", "LLM fallback вместо ручного разбора", "контроль стоимости через threshold"], 16)], fill=WHITE, rounded=True)
    box(s, 8.9, 1.45, 3.5, 3.8, [("Дополнительный агент", 23, NAVY, True), *bullets(["Financial Coach Agent", "персональные рекомендации", "tools анализа расходов", "демонстрация расширяемости"], 16)], fill=WHITE, rounded=True)
    box(s, 0.95, 5.95, 11.0, 0.55, [("Прототип нужен не как банковский продукт, а как проверка архитектурного паттерна внедрения AI-агента.", 18, NAVY, True)], PP_ALIGN.CENTER)

    # 10
    s = prs.slides.add_slide(blank)
    add_bg(s)
    title(s, "Архитектура прототипа")
    components = [
        (0.75, 1.35, "Core Service", "оркестрация, REST, БД"),
        (3.55, 1.35, "Classifier", "быстрая классификация"),
        (6.35, 1.35, "LLM Classifier", "основной AI-агент оптимизации"),
        (9.15, 1.35, "Financial Coach", "дополнительный AI-агент"),
        (2.0, 3.8, "Kafka", "topics и асинхронные события"),
        (7.2, 3.8, "PostgreSQL", "транзакции, рекомендации, JSONB"),
    ]
    for x, y, head, body in components:
        card_w = 2.45 if y < 3 else 3.4
        fill = LIGHT_BLUE if "LLM" in head or "Coach" in head else WHITE
        box(s, x, y, card_w, 1.0, [(head, 18, NAVY, True), (body, 13, MUTED, False)], PP_ALIGN.CENTER, fill, GRID, True)
    arrow(s, 3.2, 1.85, 3.5, 1.85)
    arrow(s, 6.0, 1.85, 6.3, 1.85)
    arrow(s, 8.8, 1.85, 9.1, 1.85)
    box(s, 0.9, 6.05, 11.1, 0.55, [("Ключевой принцип: LLM-компоненты изолированы как event-driven workers и не блокируют быстрый путь обработки.", 18, NAVY, True)], PP_ALIGN.CENTER)

    # 11
    s = prs.slides.add_slide(blank)
    add_bg(s)
    title(s, "Основная оптимизация: Transaction Classifier Agent")
    box(s, 0.65, 1.45, 2.0, 1.05, [("1", 24, WHITE, True), ("сырая транзакция", 14, WHITE, False)], PP_ALIGN.CENTER, NAVY, NAVY, True)
    arrow(s, 2.75, 1.98, 3.25, 1.98)
    box(s, 3.35, 1.45, 2.0, 1.05, [("2", 24, WHITE, True), ("быстрый классификатор", 14, WHITE, False)], PP_ALIGN.CENTER, NAVY, NAVY, True)
    arrow(s, 5.45, 1.98, 5.95, 1.98)
    box(s, 6.05, 1.45, 2.0, 1.05, [("3", 24, WHITE, True), ("confidence gate", 14, WHITE, False)], PP_ALIGN.CENTER, NAVY, NAVY, True)
    arrow(s, 8.15, 1.98, 8.65, 1.98)
    box(s, 8.75, 1.45, 2.0, 1.05, [("4", 24, WHITE, True), ("LLM fallback", 14, WHITE, False)], PP_ALIGN.CENTER, BLUE, BLUE, True)
    arrow(s, 10.85, 1.98, 11.35, 1.98)
    box(s, 11.45, 1.45, 1.35, 1.05, [("5", 24, WHITE, True), ("результат", 13, WHITE, False)], PP_ALIGN.CENTER, GREEN, GREEN, True)
    box(s, 0.85, 3.35, 5.45, 1.8, [("Почему это агент", 22, NAVY, True), *bullets(["получает событие из процесса", "использует контекст транзакции и историю", "принимает решение в ограниченной предметной области", "возвращает результат в Kafka/Core"], 15)], fill=LIGHT_BLUE, rounded=True)
    box(s, 6.8, 3.35, 5.45, 1.8, [("Что оптимизируется", 22, NAVY, True), *bullets(["меньше неопределённых транзакций", "меньше ручной проверки", "основной поток остаётся быстрым", "LLM вызывается только при низкой уверенности"], 15)], fill=WHITE, rounded=True)
    box(s, 0.95, 6.05, 11.1, 0.55, [("Это центральный AI-агент диплома: он улучшает существующий transaction flow, а не просто добавляет новую функцию.", 18, BLUE, True)], PP_ALIGN.CENTER)

    # 12
    s = prs.slides.add_slide(blank)
    add_bg(s)
    title(s, "Дополнительный сценарий: Financial Coach Agent")
    box(s, 0.7, 1.42, 2.4, 1.05, [("Вход", 22, NAVY, True), ("requestId, userId, период, вопрос", 14, MUTED, False)], PP_ALIGN.CENTER, WHITE, GRID, True)
    arrow(s, 3.15, 1.95, 3.7, 1.95)
    box(s, 3.85, 1.42, 2.4, 1.05, [("Tools", 22, NAVY, True), ("SQL-агрегаты расходов", 14, MUTED, False)], PP_ALIGN.CENTER, WHITE, GRID, True)
    arrow(s, 6.3, 1.95, 6.85, 1.95)
    box(s, 7.0, 1.42, 2.4, 1.05, [("LLM", 22, NAVY, True), ("совет на основе фактов", 14, MUTED, False)], PP_ALIGN.CENTER, WHITE, GRID, True)
    arrow(s, 9.45, 1.95, 10.0, 1.95)
    box(s, 10.15, 1.42, 2.4, 1.05, [("Выход", 22, NAVY, True), ("DB + coach-responses", 14, MUTED, False)], PP_ALIGN.CENTER, WHITE, GRID, True)
    box(s, 0.85, 3.45, 5.15, 1.75, [("Роль в дипломе", 22, NAVY, True), *bullets(["второй AI-агент в той же архитектуре", "демонстрация tool-based подхода", "дополнительная ценность после классификации"], 15)], fill=LIGHT_BLUE, rounded=True)
    box(s, 6.55, 3.45, 5.15, 1.75, [("Почему это не главный оптимизатор", 22, NAVY, True), *bullets(["он расширяет функциональность", "а не решает исходную проблему классификации", "поэтому основной фокус остаётся на LLM Classifier"], 15)], fill=WHITE, rounded=True)
    box(s, 0.95, 6.05, 11.1, 0.55, [("Financial Coach Agent показывает, что выбранный паттерн можно масштабировать на новые агентные сценарии.", 18, NAVY, True)], PP_ALIGN.CENTER)

    # 13
    s = prs.slides.add_slide(blank)
    add_bg(s)
    title(s, "Ожидаемый эффект оптимизации: до и после", "демонстрационные метрики MVP на синтетическом наборе")
    metric_bar(s, 0.85, 1.45, "Правильно классифицированные транзакции", 85, 95)
    metric_bar(s, 0.85, 3.0, "Транзакции UNDEFINED", 8, 1)
    metric_bar(s, 0.85, 4.55, "Ручная проверка", 5, 0.5)
    box(s, 7.2, 1.45, 4.65, 3.15, [("Интерпретация", 23, NAVY, True), *bullets(["baseline: только rule/ML слой", "после: rule/ML + LLM fallback", "LLM вызывается не для всех операций, а для low-confidence сегмента", "результат требует проверки на размеченном наборе"], 15)], fill=WHITE, rounded=True)
    box(s, 7.2, 5.05, 4.65, 0.95, [("Cost control", 18, NAVY, True), ("целевая доля LLM-вызовов: 5-10% потока", 17, BLUE, True)], PP_ALIGN.CENTER, LIGHT_BLUE, GRID, True)
    box(s, 0.9, 6.35, 11.1, 0.38, [("Цифры используются как план эксперимента и целевые ориентиры; в тексте диплома их нужно подтвердить прогоном на тестовом наборе.", 14, RED, True)], PP_ALIGN.CENTER)

    # 14
    s = prs.slides.add_slide(blank)
    add_bg(s)
    title(s, "Методика экспериментальной оценки")
    metrics = [
        ("Качество", "accuracy, доля UNDEFINED, manual review"),
        ("Операции", "p95 latency, throughput, error rate"),
        ("Экономика", "стоимость LLM-вызовов, cost per 1M tx"),
        ("Продукт", "полезность рекомендаций, feedback"),
    ]
    for x, (head, body) in zip([0.8, 3.9, 7.0, 10.1], metrics):
        box(s, x, 1.55, 2.55, 1.55, [(head, 22, NAVY, True), (body, 15, MUTED, False)], PP_ALIGN.CENTER, WHITE, GRID, True)
    box(s, 1.0, 4.0, 10.9, 1.1, [("Эксперимент", 24, NAVY, True), ("синтетический набор транзакций → baseline rule-only → hybrid rule/ML + LLM fallback → сравнение метрик", 20, BLUE, True)], PP_ALIGN.CENTER, LIGHT_BLUE, GRID, True)
    box(s, 1.0, 5.85, 10.9, 0.65, [("Важно: результаты нужно явно отделять от production-метрик, так как бизнес-кейс синтетический.", 18, RED, True)], PP_ALIGN.CENTER)

    # 15
    s = prs.slides.add_slide(blank)
    add_bg(s)
    title(s, "Выводы и ожидаемый вклад")
    box(s, 0.75, 1.45, 3.65, 3.95, [("Научно-практический вклад", 21, NAVY, True), *bullets(["систематизация AI-внедрений", "критерии сравнения технологий", "метрики оценки AI-агента"], 16)], fill=WHITE, rounded=True)
    box(s, 4.85, 1.45, 3.65, 3.95, [("Прототип", 21, NAVY, True), *bullets(["Spring/Kotlin микросервисы", "Kafka event-driven pipeline", "Transaction Classifier Agent", "Financial Coach Agent"], 16)], fill=WHITE, rounded=True)
    box(s, 8.95, 1.45, 3.65, 3.95, [("Переносимость", 21, NAVY, True), *bullets(["страхование", "кредитование", "support", "документооборот"], 16)], fill=WHITE, rounded=True)
    box(s, 0.85, 6.0, 11.25, 0.65, [("Финальная идея: LLM даёт максимум пользы не как замена всей системы, а как контролируемый агент внутри измеримого бизнес-процесса.", 18, NAVY, True)], PP_ALIGN.CENTER)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    build_slides(prs)
    prs.save(PPTX_PATH)
    print(PPTX_PATH)


if __name__ == "__main__":
    main()
